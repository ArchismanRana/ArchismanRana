from tkinter import *
import customtkinter as ctk
from tkinter import messagebox
import win32com
from PIL import Image
from tkinter import filedialog
import pathlib
from pdf2docx import Converter
from win32com.client import Dispatch
import os
import sys
import PyPDF2
import img2pdf
import pdfplumber
import pandas as pd  # pip install openpyxl
from pdf2image import convert_from_path
from pptx import Presentation
from pptx.util import Inches
from docx.shared import Inches
from docx2pdf import convert


def resource_path(relative_path):
    """ Get the absolute path to the resource, works for both development and PyInstaller """
    try:
        base_path = sys._MEIPASS
    except Exception:
        base_path = os.path.abspath(".")

    return os.path.join(base_path, relative_path)


root = ctk.CTk()
root.title("PDF converter and merging PDFs")
ctk.set_appearance_mode("System")
root.iconbitmap('photos/APP LOGO.ico')
root.configure(fg_color="#7b2929")

root.geometry('743x550')  # Order is WidthXHeight
root.minsize(743, 550)  # Order is Width,Height
root.maxsize(743, 550)
root.configure(bg='brown')

# Animation
my_y = -70


# Animation Functions
def down():
    global my_y
    if my_y <= 10:
        my_y += 10
        Text_1.place(x=0, y=my_y)
        root.after(40, down)


Frame1 = ctk.CTkFrame(root, corner_radius=10)
Frame1.pack(padx=20, pady=20, fill="both", expand=True)
Frame2 = ctk.CTkFrame(master=Frame1, corner_radius=10, width=200, height=100, fg_color="white")
Frame2.pack(padx=10, pady=10, fill="both", expand=True)
Text_1 = ctk.CTkLabel(Frame2, text="Welcome to PDF converter and Merging PDF", text_color="dark blue",
                      fg_color="#eeeeee", font=("Microsoft New Tai Lue", 30, "bold"),
                      corner_radius=8, padx=15, pady=15)
Text_1.place(x=0, y=my_y)
root.after(400, down)

img1 = ctk.CTkImage(light_image=Image.open(resource_path('photos/WordToPDF.png')), size=(166, 90))

img2 = ctk.CTkImage(light_image=Image.open(resource_path('photos/pdftoword.png')), size=(166, 90))

img3 = ctk.CTkImage(light_image=Image.open(resource_path('photos/Exit.png')), size=(60, 55))

img4 = ctk.CTkImage(light_image=Image.open(resource_path('photos/ExcelToPDF.png')), size=(166, 90))

img5 = ctk.CTkImage(light_image=Image.open(resource_path('photos/pdftoexcel.png')), size=(166, 90))

img7 = ctk.CTkImage(light_image=Image.open(resource_path('photos/PptToPdf.png')), size=(166, 90))

img8 = ctk.CTkImage(light_image=Image.open(resource_path('photos/Pdf_to_PPT.png')), size=(166, 90))

img9 = ctk.CTkImage(light_image=Image.open(resource_path('photos/JPGToPDF.png')), size=(166, 90))

img10 = ctk.CTkImage(light_image=Image.open(resource_path('photos/PdfToPng.png')), size=(166, 90))

img11 = ctk.CTkImage(light_image=Image.open(resource_path('photos/Insert_Page.png')), size=(166, 100))

img12 = ctk.CTkImage(light_image=Image.open(resource_path("photos/Remove_Page.png")), size=(100, 100))


def convert_doc_to_docx(doc_file_path, docx_file_path):
    word = win32com.client.Dispatch("Word.Application")
    word.Visible = False

    # Ensure the file path is absolute and formatted correctly
    doc_file_path_n = os.path.abspath(doc_file_path)
    docx_file_path_n = os.path.abspath(docx_file_path)

    if os.path.exists(doc_file_path):
        doc = word.Documents.Open(doc_file_path)
        doc.SaveAs(docx_file_path, FileFormat=16)  # 16 corresponds to the DOCX format
        doc.Close()
        word.Quit()
        messagebox.showinfo("Information", f"Successfully converted {doc_file_path_n} to {docx_file_path_n}")
    else:
        messagebox.showerror("Error", f"File not found: {doc_file_path}")


def w_to_p():
    New_Win = ctk.CTkToplevel(root)  # TopLevel object which will be treated as a new window
    New_Win.title("Word to pdf")
    New_Win.geometry('800x350')  # Order is WidthXHeight
    New_Win.minsize(800, 400)
    New_Win.maxsize(1000, 600)
    New_Win.configure(fg_color='#b2baef')

    heading = ctk.CTkLabel(New_Win, text="Word to PDF converter",
                           text_color="dark blue", fg_color="white",
                           font=("Berlin Sans FB", 30))
    heading.pack(fill='x', pady='14')

    Doc = ctk.CTkLabel(New_Win, text="Enter File path:", fg_color='#b2baef',
                       text_color='black', font=('Franklin Gothic Demi', 30, 'bold'))
    Doc.place(x=70, y=90)
    file_path_var = StringVar()
    En = ctk.CTkEntry(New_Win, height=25, width=600, font=("helvetica", 20),
                      corner_radius=20, fg_color="white",
                      text_color="green", textvariable=file_path_var)
    En.place(x=70, y=150)

    def open_file():
        file = filedialog.askopenfilename(filetypes=[("Word files", "*.docx *.doc"), ("All files", "*.*")])
        if file:
            file_path_var.set(file)

    def convert_file():
        file_path = file_path_var.get()
        if pathlib.Path(file_path).suffix == '.docx':
            convert(file_path)
            messagebox.showinfo("Information", "File converted successfully!")
            New_Win.destroy()
        elif pathlib.Path(file_path).suffix == '.doc':
            docx_file_path = file_path.replace('.doc', '.docx')
            convert_doc_to_docx(file_path, docx_file_path)
            convert(docx_file_path)
            messagebox.showinfo("Information", "File converted successfully!")
            New_Win.destroy()
        else:
            messagebox.showerror("Error", "Please select a file with .docx extension")
            New_Win.destroy()

    browse_button = ctk.CTkButton(New_Win, text="Browse", fg_color="white", text_color='green',
                                  font=('Verdana', 25, 'bold'), width=12, height=12,
                                  hover_color="#f9cb9c", border_width=2, border_color="grey", command=open_file)
    browse_button.place(x=80, y=210)

    convert_button = ctk.CTkButton(New_Win, text="Convert", fg_color="white", text_color='red',
                                   font=('Verdana', 25, 'bold'), width=12, height=12,
                                   hover_color="#f9cb9c", border_width=2, border_color="grey", command=convert_file)
    convert_button.place(x=80, y=265)


def convert_excel_to_pdf(excel_file_path, pdf_file_path):
    excel = win32com.client.Dispatch("Excel.Application")
    excel.Visible = False

    try:
        # Ensure the file path is absolute and formatted correctly
        excel_file_path = os.path.abspath(excel_file_path)
        pdf_file_path = os.path.abspath(pdf_file_path)

        if os.path.exists(excel_file_path):
            wb = excel.Workbooks.Open(excel_file_path)
            wb.ExportAsFixedFormat(0, pdf_file_path)  # 0 corresponds to PDF format
            wb.Close(False)
            excel.Quit()
            messagebox.showinfo("Information", f"Successfully converted {excel_file_path} to {pdf_file_path}")
        else:
            messagebox.showerror("Error", f"File not found: {excel_file_path}")
    except Exception as e:
        excel.Quit()
        messagebox.showerror("Error", f"Failed to convert {excel_file_path} to PDF: {e}")


def e_to_p():
    New_Win = ctk.CTkToplevel(root)  # TopLevel object which will be treated as a new window
    New_Win.title("Excel to PDF")
    New_Win.geometry('800x350')  # Order is WidthXHeight
    New_Win.minsize(800, 400)
    New_Win.maxsize(1000, 600)
    New_Win.configure(fg_color='#b2baef')

    heading = ctk.CTkLabel(New_Win, text="Excel to PDF Converter",
                           text_color="dark blue", fg_color="white",
                           font=("Berlin Sans FB", 30))
    heading.pack(fill='x', pady='14')

    Doc = ctk.CTkLabel(New_Win, text="Enter File path:", fg_color='#b2baef',
                       text_color='black', font=('Franklin Gothic Demi', 30, 'bold'))
    Doc.place(x=70, y=90)
    file_path_var = StringVar()
    En = ctk.CTkEntry(New_Win, height=25, width=600, font=("helvetica", 20),
                      corner_radius=20, fg_color="white",
                      text_color="green", textvariable=file_path_var)
    En.place(x=70, y=150)

    def open_file():
        file = filedialog.askopenfilename(filetypes=[("Excel files", "*.xlsx *.xls"), ("All files", "*.*")])
        if file:
            file_path_var.set(file)

    def convert_file():
        file_path = file_path_var.get()
        pdf_file_path = file_path.replace('.xls', '.pdf').replace('.xlsx', '.pdf')
        convert_excel_to_pdf(file_path, pdf_file_path)

    browse_button = ctk.CTkButton(New_Win, text="Browse", fg_color="white", text_color='green',
                                  font=('Verdana', 25, 'bold'), width=12, height=12,
                                  hover_color="#f9cb9c", border_width=2, border_color="grey", command=open_file)
    browse_button.place(x=80, y=210)

    convert_button = ctk.CTkButton(New_Win, text="Convert", fg_color="white", text_color='red',
                                   font=('Verdana', 25, 'bold'), width=12, height=12,
                                   hover_color="#f9cb9c", border_width=2, border_color="grey", command=convert_file)
    convert_button.place(x=80, y=265)


def add_poppler_to_path():
    # Specify the path to the Poppler bin directory
    poppler_path = r'C:\poppler\poppler-24.02.0\Library\bin'
    # Add Poppler bin directory to the PATH environment variable
    os.environ['PATH'] += os.pathsep + poppler_path


# Function to convert PPT to PDF
def convert_ppt_to_pdf(ppt_path, pdf_path):
    powerpoint = win32com.client.Dispatch("PowerPoint.Application")
    powerpoint.Visible = True

    try:
        # Ensure the file path is absolute and formatted correctly
        ppt_file_path = os.path.abspath(ppt_path)
        pdf_file_path = os.path.abspath(pdf_path)

        if os.path.exists(ppt_file_path):
            presentation = powerpoint.Presentations.Open(ppt_file_path)
            presentation.SaveAs(pdf_file_path, 32)  # 32 corresponds to the PDF format
            presentation.Close()
            powerpoint.Quit()
            messagebox.showinfo("Information", f"Successfully converted {ppt_file_path} to {pdf_file_path}")
        else:
            messagebox.showerror("Error", f"File not found: {ppt_file_path}")
    except Exception as e:
        powerpoint.Quit()
        messagebox.showerror("Error", f"Failed to convert ppt file to PDF: {e}")


def p_to_p():
    New_Win = ctk.CTkToplevel(root)  # TopLevel object which will be treated as a new window
    New_Win.title("PPT to PDF")
    New_Win.geometry('800x350')  # Order is WidthXHeight
    New_Win.minsize(800, 400)
    New_Win.maxsize(1000, 600)
    New_Win.configure(fg_color='#b2baef')

    heading = ctk.CTkLabel(New_Win, text="PPT to PDF Converter",
                           text_color="dark blue", fg_color="white",
                           font=("Berlin Sans FB", 30))
    heading.pack(fill='x', pady='14')

    ctk.CTkLabel(New_Win, text="Enter File path:", fg_color='#b2baef',
                 text_color='black', font=('Franklin Gothic Demi', 30, 'bold')).place(x=70, y=90)
    file_path_var = StringVar()
    ctk.CTkEntry(New_Win, height=25, width=600, font=("helvetica", 20),
                 corner_radius=20, fg_color="white",
                 text_color="green", textvariable=file_path_var).place(x=70, y=150)

    def open_file():
        file = filedialog.askopenfilename(filetypes=[("PowerPoint files", "*.ppt;*.pptx")])
        if file:
            file_path_var.set(file)

    def convert_file():
        file_path = file_path_var.get()
        if pathlib.Path(file_path).suffix in ['.ppt', '.pptx']:
            pdf_path = file_path.replace('.ppt', '.pdf').replace('.pptx', '.pdf')
            convert_ppt_to_pdf(file_path, pdf_path)
            New_Win.destroy()
        else:
            messagebox.showerror("Error", "Please select a file with .ppt or .pptx extension")

    ctk.CTkButton(New_Win, text="Browse", fg_color="white", text_color='green',
                  font=('Verdana', 25, 'bold'), width=12, height=12,
                  hover_color="#f9cb9c", border_width=2, border_color="grey",
                  command=open_file).place(x=80, y=210)

    ctk.CTkButton(New_Win, text="Convert", fg_color="white", text_color='red',
                  font=('Verdana', 25, 'bold'), width=12, height=12,
                  hover_color="#f9cb9c", border_width=2, border_color="grey",
                  command=convert_file).place(x=80, y=265)


def i_to_p():
    New_Win = ctk.CTkToplevel(root)  # TopLevel object which will be treated as a new window
    New_Win.title("Image to pdf")
    New_Win.geometry('800x350')  # Order is WidthXHeight
    New_Win.minsize(800, 400)
    New_Win.maxsize(1000, 600)
    New_Win.configure(fg_color='#b2baef')

    heading = ctk.CTkLabel(New_Win, text="Image to PDF converter",
                           text_color="dark blue", fg_color="white",
                           font=("Berlin Sans FB", 30))
    heading.pack(fill='x', pady='14')

    Doc = ctk.CTkLabel(New_Win, text="Enter File path:", fg_color='#b2baef',
                       text_color='black', font=('Franklin Gothic Demi', 30, 'bold'))
    Doc.place(x=70, y=90)
    file_path_var = StringVar()
    En = ctk.CTkEntry(New_Win, height=25, width=600, font=("helvetica", 20),
                      corner_radius=20, fg_color="white",
                      text_color="green", textvariable=file_path_var)
    En.place(x=70, y=150)

    def open_file():
        file = filedialog.askopenfilename(filetypes=[("Image Files", "*.png;*.jpg;*.jpeg;*.bmp;*.tiff")])
        if file:
            file_path_var.set(file)

    def convert_file():
        file_path = file_path_var.get()
        if pathlib.Path(file_path).suffix in ['.png', '.jpg', '.jpeg', '.bmp', '.tiff']:
            save_as = filedialog.asksaveasfilename(defaultextension=".pdf", filetypes=[("PDF Files", "*.pdf")])
            if save_as:
                try:
                    with open(save_as, "wb") as f:
                        f.write(img2pdf.convert(file_path))
                    messagebox.showinfo("Information", "File converted successfully!")
                    New_Win.destroy()
                except Exception as e:
                    messagebox.showerror("Error", f"Failed to convert file: {e}")
                    New_Win.destroy()
        else:
            messagebox.showerror("Error", "Please select a valid image file")
            New_Win.destroy()

    browse_button = ctk.CTkButton(New_Win, text="Browse", fg_color="white", text_color='green',
                                  font=('Verdana', 25, 'bold'), width=12, height=12,
                                  hover_color="#f9cb9c", border_width=2, border_color="grey", command=open_file)
    browse_button.place(x=80, y=210)

    convert_button = ctk.CTkButton(New_Win, text="Convert", fg_color="white", text_color='red',
                                   font=('Verdana', 25, 'bold'), width=12, height=12,
                                   hover_color="#f9cb9c", border_width=2, border_color="grey", command=convert_file)
    convert_button.place(x=80, y=265)


def merge_pdfs(filepaths, output_path):
    pdf_writer = PyPDF2.PdfWriter()  # Blank PDF
    for i in filepaths:
        if i:
            try:
                pdf_reader = PyPDF2.PdfReader(i)
                for page_num in range(len(pdf_reader.pages)):
                    page = pdf_reader.pages[page_num]
                    pdf_writer.add_page(page)  # PDFs added to blank pdf page by page
            except FileNotFoundError:
                return
            except Exception as e:
                messagebox.showerror("Error", f"An error occurred: {e}")
                return

        with open(output_path, 'wb') as out:
            pdf_writer.write(out)


def merging():
    New_Win = ctk.CTkToplevel(root)  # TopLevel object which will be treated as a new window
    New_Win.title("Merge pdf's")
    New_Win.geometry('860x500')  # Order is WidthXHeight
    New_Win.minsize(800, 400)
    New_Win.maxsize(900, 600)
    New_Win.configure(fg_color='#b2baef')

    my_frame = ctk.CTkScrollableFrame(New_Win, width=800, height=300, fg_color="white",
                                      scrollbar_button_color="orange", scrollbar_button_hover_color="white",
                                      scrollbar_fg_color="black",
                                      label_text="Merge Four PDFs into a single PDF", label_fg_color="white",
                                      label_text_color="dark blue", label_font=("Agency FB", 40, "bold"),
                                      label_anchor="center", border_width=5, border_color="Black",
                                      corner_radius=4)
    my_frame.pack(padx=20, pady=20, fill="both", expand=True)

    def open_1file():
        files = filedialog.askopenfilename(filetypes=[("PDF files", "*.pdf")])
        if files:
            file_path_var_1.set(files)

    def open_2file():
        files = filedialog.askopenfilename(filetypes=[("PDF files", "*.pdf")])
        if files:
            file_path_var_2.set(files)

    def open_3file():
        files = filedialog.askopenfilename(filetypes=[("PDF files", "*.pdf")])
        if files:
            file_path_var_3.set(files)

    def open_4file():
        files = filedialog.askopenfilename(filetypes=[("PDF files", "*.pdf")])
        if files:
            file_path_var_4.set(files)

    Doc_1 = ctk.CTkLabel(my_frame, text="♦Enter first File path:", fg_color='white',
                         text_color='black', font=('Franklin Gothic Demi', 30, 'bold'))
    Doc_1.pack(padx=55, pady=10, anchor='w')
    file_path_var_1 = StringVar()
    En1 = ctk.CTkEntry(my_frame,
                       height=25, width=700, font=("helvetica", 20),
                       corner_radius=20, fg_color="white",
                       text_color="green", textvariable=file_path_var_1)
    En1.pack(padx=5, pady=10)
    Doc_1entry = ctk.CTkButton(my_frame, text="Browse File", fg_color="white", text_color='green',
                               font=('Verdana', 20, 'bold'), width=12, height=12, command=open_1file,
                               border_width=3, border_color="grey", hover_color="#f9cb9c")
    Doc_1entry.pack(padx=55, pady=10, anchor='w')

    Doc_2 = ctk.CTkLabel(my_frame, text="♦Enter second File path:", fg_color='white',
                         text_color='black', font=('Franklin Gothic Demi', 30, 'bold'))
    Doc_2.pack(padx=55, pady=10, anchor='w')
    file_path_var_2 = StringVar()
    En2 = ctk.CTkEntry(my_frame,
                       height=25, width=700, font=("helvetica", 20),
                       corner_radius=20, fg_color="white",
                       text_color="green", textvariable=file_path_var_2)
    En2.pack(padx=5, pady=10)
    Doc_2entry = ctk.CTkButton(my_frame, text="Browse File", fg_color="white", text_color='green',
                               font=('Verdana', 20, 'bold'), width=12, height=12, command=open_2file,
                               border_width=3, border_color="grey", hover_color="#f9cb9c")
    Doc_2entry.pack(padx=55, pady=10, anchor='w')

    Doc_3 = ctk.CTkLabel(my_frame, text="♦Enter third File path:", fg_color='white',
                         text_color='black', font=('Franklin Gothic Demi', 30, 'bold'))
    Doc_3.pack(padx=55, pady=10, anchor='w')
    file_path_var_3 = StringVar()
    En3 = ctk.CTkEntry(my_frame,
                       height=25, width=700, font=("helvetica", 20),
                       corner_radius=20, fg_color="white",
                       text_color="green", textvariable=file_path_var_3)
    En3.pack(padx=5, pady=10)
    Doc_3entry = ctk.CTkButton(my_frame, text="Browse File", fg_color="white", text_color='green',
                               font=('Verdana', 20, 'bold'), width=12, height=12, command=open_3file,
                               border_width=3, border_color="grey", hover_color="#f9cb9c")
    Doc_3entry.pack(padx=55, pady=10, anchor='w')

    Doc_4 = ctk.CTkLabel(my_frame, text="♦Enter fourth File path:", fg_color='white',
                         text_color='black', font=('Franklin Gothic Demi', 30, 'bold'))
    Doc_4.pack(padx=55, pady=10, anchor='w')
    file_path_var_4 = StringVar()
    En4 = ctk.CTkEntry(my_frame,
                       height=25, width=700, font=("helvetica", 20),
                       corner_radius=20, fg_color="white",
                       text_color="green", textvariable=file_path_var_4)
    En4.pack(padx=5, pady=10)
    Doc_4entry = ctk.CTkButton(my_frame, text="Browse File", fg_color="white", text_color='green',
                               font=('Verdana', 20, 'bold'), width=12, height=12, command=open_4file,
                               border_width=3, border_color="grey", hover_color="#f9cb9c")
    Doc_4entry.pack(padx=55, pady=20, anchor='w')

    def merge_files():
        f1 = file_path_var_1.get()
        f2 = file_path_var_2.get()
        f3 = file_path_var_3.get()
        f4 = file_path_var_4.get()
        files = [f1, f2, f3, f4]
        valid_files = [f for f in files if f]  # Filter out empty file paths
        if len(valid_files) < 2:
            messagebox.showerror("Error", "Select at least two PDF files to merge.")
            return
        output_path = filedialog.asksaveasfilename(defaultextension=".pdf", filetypes=[("PDF files", "*.pdf")])
        if output_path:
            merge_pdfs(valid_files, output_path)
            messagebox.showinfo("Information", "PDFs merged successfully!")
            New_Win.destroy()

    merge_button = ctk.CTkButton(my_frame, text="Merge PDFs'  ", fg_color="white", text_color='red',
                                 font=('Verdana', 20, 'bold'), width=12, height=12, command=merge_files,
                                 border_width=3, border_color="grey", hover_color="#f9cb9c")
    merge_button.place(x=540, y=610)


def close():
    x = messagebox.askquestion("askquestion", "Are you sure?")
    if x == 'yes':
        root.destroy()
    else:
        return None


def doc_to_pdf():
    New_Win = ctk.CTkToplevel(root, fg_color='#b2baef')  # TopLevel object which will be treated as a new window
    New_Win.title("Docx to Pdf")
    New_Win.geometry('800x550')  # Order is WidthXHeight
    New_Win.minsize(800, 500)
    New_Win.maxsize(1000, 600)

    heading = ctk.CTkLabel(New_Win, text="Convert Doc to PDF",
                           text_color="dark blue", fg_color="white", height=50,
                           font=("Microsoft New Tai Lue", 30, "bold"))
    heading.pack(fill='x', pady='15')

    o = ctk.CTkButton(New_Win, fg_color="white", text="Word to PDF", text_color="dark blue",
                      command=w_to_p, font=("Comic Sans MS", 20, "bold"), hover_color="#f9cb9c",
                      image=img1, compound="top", border_width=5, border_color="#291856")
    o.place(x=100, y=105)

    p = ctk.CTkButton(New_Win, fg_color="white", text="Excel to PDF", text_color="dark blue",
                      command=e_to_p, font=("Comic Sans MS", 20, "bold"), hover_color="#f9cb9c",
                      image=img4, compound="top", border_width=5, border_color="#291856")
    p.place(x=500, y=105)

    q = ctk.CTkButton(New_Win, fg_color="white", text="PPT to PDF", text_color="dark blue",
                      command=p_to_p, font=("Comic Sans MS", 20, "bold"), hover_color="#f9cb9c",
                      image=img7, compound="top", border_width=5, border_color="#291856")
    q.place(x=100, y=305)

    r = ctk.CTkButton(New_Win, fg_color="white", text="Image to PDF", text_color="dark blue",
                      command=i_to_p, font=("Comic Sans MS", 20, "bold"), hover_color="#f9cb9c",
                      image=img9, compound="top", border_width=5, border_color="#291856")
    r.place(x=500, y=305)


def p_to_w():
    New_Win = ctk.CTkToplevel(root)  # TopLevel object which will be treated as a new window
    New_Win.title("pdf to Word")
    New_Win.geometry('800x350')  # Order is WidthXHeight
    New_Win.minsize(800, 400)
    New_Win.maxsize(1000, 600)
    New_Win.configure(fg_color='#b2baef')

    heading = ctk.CTkLabel(New_Win, text="PDF to Word converter",
                           text_color="dark blue", fg_color="white",
                           font=("Berlin Sans FB", 30))
    heading.pack(fill='x', pady='14')

    Doc = ctk.CTkLabel(New_Win, text="Enter File path:", fg_color='#b2baef',
                       text_color='black', font=('Franklin Gothic Demi', 30, 'bold'))
    Doc.place(x=70, y=90)
    file_path_var = StringVar()
    En = ctk.CTkEntry(New_Win, height=25, width=600, font=("helvetica", 20),
                      corner_radius=20, fg_color="white",
                      text_color="green", textvariable=file_path_var)
    En.place(x=70, y=150)

    def open_file():
        file = filedialog.askopenfilename(filetypes=[("Pdf files", "*.pdf")])
        if file:
            file_path_var.set(file)

    def convert_file():
        file_path = file_path_var.get()
        if pathlib.Path(file_path).suffix == '.pdf':
            docx_path = file_path.replace(".pdf", ".docx")
            cv = Converter(file_path)
            cv.convert(docx_path)
            cv.close()
            messagebox.showinfo("Information", "File converted successfully!")
            New_Win.destroy()
        else:
            messagebox.showerror("Error", "Failed to convert file. Please select a file with .docx extension")
            New_Win.destroy()

    browse_button = ctk.CTkButton(New_Win, text="Browse", fg_color="white", text_color='green',
                                  font=('Verdana', 25, 'bold'), width=12, height=12,
                                  hover_color="#f9cb9c", border_width=2, border_color="grey", command=open_file)
    browse_button.place(x=80, y=210)

    convert_button = ctk.CTkButton(New_Win, text="Convert", fg_color="white", text_color='red',
                                   font=('Verdana', 25, 'bold'), width=12, height=12,
                                   hover_color="#f9cb9c", border_width=2, border_color="grey", command=convert_file)
    convert_button.place(x=80, y=265)


def p_to_e():
    New_Win = ctk.CTkToplevel(root)  # TopLevel object which will be treated as a new window
    New_Win.title("PDF to Excel")
    New_Win.geometry('800x350')  # Order is WidthXHeight
    New_Win.minsize(800, 400)
    New_Win.maxsize(1000, 600)
    New_Win.configure(fg_color='#b2baef')

    heading = ctk.CTkLabel(New_Win, text="PDF to Excel converter",
                           text_color="dark blue", fg_color="white",
                           font=("Berlin Sans FB", 30))
    heading.pack(fill='x', pady='14')

    Doc = ctk.CTkLabel(New_Win, text="Enter File path:", fg_color='#b2baef',
                       text_color='black', font=('Franklin Gothic Demi', 30, 'bold'))
    Doc.place(x=70, y=90)
    file_path_var = StringVar()
    En = ctk.CTkEntry(New_Win, width=600, font=("helvetica", 20),
                      corner_radius=20, fg_color="white",
                      text_color="green", textvariable=file_path_var)
    En.place(x=70, y=150)

    def open_file():
        file = filedialog.askopenfilename(filetypes=[("Pdf files", "*.pdf")])
        if file:
            file_path_var.set(file)

    def convert_file():
        file_path = file_path_var.get()
        if pathlib.Path(file_path).suffix == '.pdf':
            save_as = filedialog.asksaveasfilename(defaultextension=".xlsx", filetypes=[("Excel Files", "*.xlsx")])
            if save_as:
                try:
                    with pdfplumber.open(file_path) as pdf:
                        all_tables = []
                        for page in pdf.pages:
                            tables = page.extract_tables()
                            for table in tables:
                                df = pd.DataFrame(table)
                                all_tables.append(df)

                        if all_tables:
                            result = pd.concat(all_tables, ignore_index=True)
                            result.to_excel(save_as, index=False)
                            messagebox.showinfo("Information", "File converted successfully!")
                            New_Win.destroy()
                        else:
                            messagebox.showerror("Error", "No tables found in the PDF")
                except Exception as e:
                    messagebox.showerror("Error", f"Failed to convert file: {e}")
                    New_Win.destroy()
        else:
            messagebox.showerror("Error", "Please select a valid PDF file")
            New_Win.destroy()

    browse_button = ctk.CTkButton(New_Win, text="Browse", fg_color="white", text_color='green',
                                  font=('Verdana', 25, 'bold'), width=12, height=12,
                                  hover_color="#f9cb9c", border_width=2, border_color="grey", command=open_file)
    browse_button.place(x=80, y=210)

    convert_button = ctk.CTkButton(New_Win, text="Convert", fg_color="white", text_color='red',
                                   font=('Verdana', 25, 'bold'), width=12, height=12,
                                   hover_color="#f9cb9c", border_width=2, border_color="grey", command=convert_file)
    convert_button.place(x=80, y=265)


def pdf_to_ppt(pdf_path, output_path):
    images = convert_from_path(pdf_path)
    presentation = Presentation()

    for image in images:
        slide_layout = presentation.slide_layouts[6]
        slide = presentation.slides.add_slide(slide_layout)
        imag_path = os.path.join(output_path, "slide_image.jpg")
        image.save(imag_path, 'JPEG')
        left = top = Inches(0)
        slide.shapes.add_picture(imag_path, left, top, width=Inches(10), height=Inches(7.5))

    pptx_output_path = os.path.join(output_path, "output_presentation.pptx")
    presentation.save(pptx_output_path)


def pd_to_pp():
    add_poppler_to_path()
    New_Win = ctk.CTkToplevel(root)  # TopLevel object which will be treated as a new window
    New_Win.title("pdf to PPT")
    New_Win.geometry('800x350')  # Order is WidthXHeight
    New_Win.minsize(800, 400)
    New_Win.maxsize(1000, 600)
    New_Win.configure(fg_color='#b2baef')

    heading = ctk.CTkLabel(New_Win, text="PDF to PPT converter",
                           text_color="dark blue", fg_color="white",
                           font=("Berlin Sans FB", 30))
    heading.pack(fill='x', pady='15')

    Doc = ctk.CTkLabel(New_Win, text="Enter File path:", fg_color='#b2baef',
                       text_color='black', font=('Franklin Gothic Demi', 30, 'bold'))
    Doc.place(x=70, y=90)
    file_path_var = StringVar()
    En = ctk.CTkEntry(New_Win, height=25, width=600, font=("helvetica", 20),
                      corner_radius=20, fg_color="white",
                      text_color="green", textvariable=file_path_var)
    En.place(x=70, y=150)

    def open_file():
        file = filedialog.askopenfilename(filetypes=[("Pdf files", "*.pdf")])
        if file:
            file_path_var.set(file)

    def convert_file():
        file_path = file_path_var.get()
        if pathlib.Path(file_path).suffix == '.pdf':
            output_dir = filedialog.askdirectory()
            if output_dir:
                pdf_to_ppt(file_path, output_dir)
                messagebox.showinfo("Information", "File converted successfully!")
                New_Win.destroy()

        else:
            messagebox.showerror("Error", "Failed to convert file. Select a file with .pdf extension")
            New_Win.destroy()

    browse_button = ctk.CTkButton(New_Win, text="Browse", fg_color="white", text_color='green',
                                  font=('Verdana', 25, 'bold'), width=12, height=12,
                                  hover_color="#f9cb9c", border_width=2, border_color="grey", command=open_file)
    browse_button.place(x=80, y=210)

    convert_button = ctk.CTkButton(New_Win, text="Convert", fg_color="white", text_color='red',
                                   font=('Verdana', 25, 'bold'), width=12, height=12,
                                   hover_color="#f9cb9c", border_width=2, border_color="grey", command=convert_file)
    convert_button.place(x=80, y=265)


def p_to_i():
    # Function to convert PDF to images using pdf2image
    add_poppler_to_path()
    New_Win = ctk.CTkToplevel(root)  # TopLevel object which will be treated as a new window
    New_Win.title("PDF to Image")
    New_Win.geometry('800x350')  # Order is WidthXHeight
    New_Win.minsize(800, 400)
    New_Win.maxsize(1000, 600)
    New_Win.configure(fg_color='#b2baef')

    heading = ctk.CTkLabel(New_Win, text="PDF to Image converter",
                           text_color="dark blue", fg_color="white",
                           font=("Berlin Sans FB", 30))
    heading.pack(fill='x', pady='15')

    Doc = ctk.CTkLabel(New_Win, text="Enter File path:", fg_color='#b2baef',
                       text_color='black', font=('Franklin Gothic Demi', 30, 'bold'))
    Doc.place(x=70, y=90)
    file_path_var = StringVar()
    En = ctk.CTkEntry(New_Win, height=25, width=600, font=("helvetica", 20),
                      corner_radius=20, fg_color="white",
                      text_color="green", textvariable=file_path_var)
    En.place(x=70, y=150)

    def open_file():
        file = filedialog.askopenfilename(filetypes=[("Pdf files", "*.pdf")])
        if file:
            file_path_var.set(file)

    def convert_file():
        file_path = file_path_var.get()
        if pathlib.Path(file_path).suffix == '.pdf':
            save_dir = filedialog.askdirectory()
            if save_dir:
                try:
                    images = convert_from_path(file_path)
                    for i, image in enumerate(images):
                        image.save(f"{save_dir}/page_{i + 1}.png", "PNG")
                    messagebox.showinfo("Information", "PDF converted to images successfully!")
                    New_Win.destroy()
                except Exception as e:
                    messagebox.showerror("Error", f"Failed to convert file: {e}")
                    New_Win.destroy()
        else:
            messagebox.showerror("Error", "Please select a valid PDF file")
            New_Win.destroy()

    browse_button = ctk.CTkButton(New_Win, text="Browse", fg_color="white", text_color='green',
                                  font=('Verdana', 25, 'bold'), width=12, height=12,
                                  hover_color="#f9cb9c", border_width=2, border_color="grey", command=open_file)
    browse_button.place(x=80, y=210)

    convert_button = ctk.CTkButton(New_Win, text="Convert", fg_color="white", text_color='red',
                                   font=('Verdana', 25, 'bold'), width=12, height=12,
                                   hover_color="#f9cb9c", border_width=2, border_color="grey", command=convert_file)
    convert_button.place(x=80, y=265)


def pdf_to_doc():
    New_Win = ctk.CTkToplevel(root)  # TopLevel object which will be treated as a new window
    New_Win.title("PDF to Docx")
    New_Win.geometry('800x550')  # Order is WidthXHeight
    New_Win.minsize(800, 500)
    New_Win.maxsize(1000, 600)
    New_Win.configure(fg_color='#b2baef')

    heading = ctk.CTkLabel(New_Win, text="Convert PDF to Doc",
                           text_color="dark blue", fg_color="white", height=50,
                           font=("Microsoft New Tai Lue", 30, "bold"))
    heading.pack(fill='x', pady='15')

    o = ctk.CTkButton(New_Win, fg_color="white", text="PDF to Word", text_color="dark blue", command=p_to_w,
                      font=("Comic Sans MS", 20, "bold"), hover_color="#f9cb9c",
                      image=img2, compound="top", border_width=5, border_color="#291856")
    o.place(x=100, y=105)

    o = ctk.CTkButton(New_Win, fg_color="white", text="PDF to Excel", text_color="dark blue", command=p_to_e,
                      font=("Comic Sans MS", 20, "bold"), hover_color="#f9cb9c",
                      image=img5, compound="top", border_width=5, border_color="#291856")
    o.place(x=500, y=105)

    o = ctk.CTkButton(New_Win, fg_color="white", text="PDF to PPT", text_color="dark blue", command=pd_to_pp,
                      font=("Comic Sans MS", 20, "bold"), hover_color="#f9cb9c",
                      image=img8, compound="top", border_width=5, border_color="#291856")
    o.place(x=100, y=305)

    o = ctk.CTkButton(New_Win, fg_color="white", text="PDF to Image", text_color="dark blue", command=p_to_i,
                      font=("Comic Sans", 20, "bold"), hover_color="#f9cb9c",
                      image=img10, compound="top", border_width=5, border_color="#291856")
    o.place(x=500, y=305)


def insert_page_into_pdf(existing_pdf_path, insert_pdf_path, position, output_path):
    existing_pdf_reader = PyPDF2.PdfReader(existing_pdf_path)
    insert_pdf_reader = PyPDF2.PdfReader(insert_pdf_path)
    pdf_writer = PyPDF2.PdfWriter()

    for page_num in range(position):
        pdf_writer.add_page(existing_pdf_reader.pages[page_num])

    for page in insert_pdf_reader.pages:
        pdf_writer.add_page(page)

    for page_num in range(position, len(existing_pdf_reader.pages)):
        pdf_writer.add_page(existing_pdf_reader.pages[page_num])

    with open(output_path, 'wb') as out_file:
        pdf_writer.write(out_file)


def insert():
    New_Win = ctk.CTkToplevel(root)  # TopLevel object which will be treated as a new window
    New_Win.title("Edit PDF")
    New_Win.geometry('750x650')  # Order is WidthXHeight
    New_Win.minsize(750, 600)
    New_Win.maxsize(1000, 600)
    New_Win.configure(fg_color='#b2baef')

    heading = ctk.CTkLabel(New_Win, text="Insert page to an existing PDF",
                           text_color="dark blue", fg_color="white", height=50,
                           font=("Microsoft New Tai Lue", 30, "bold"))
    heading.pack(fill='x', pady='14')

    file_path_var = StringVar()
    insert_path_var = StringVar()
    position_var = StringVar()

    def open_existing_file():
        file = filedialog.askopenfilename(filetypes=[("PDF files", "*.pdf")])
        if file:
            file_path_var.set(file)

    def open_insert_file():
        file = filedialog.askopenfilename(filetypes=[("PDF files", "*.pdf")])
        if file:
            insert_path_var.set(file)

    def insert_file():
        existing_pdf_path = file_path_var.get()
        insert_pdf_path = insert_path_var.get()
        position = int(position_var.get())
        output_path = filedialog.asksaveasfilename(defaultextension=".pdf", filetypes=[("PDF files", "*.pdf")])
        if output_path:
            insert_page_into_pdf(existing_pdf_path, insert_pdf_path, position, output_path)
            messagebox.showinfo("Information", "Page inserted successfully!")
            New_Win.destroy()

    ctk.CTkLabel(New_Win, text="►Enter existing PDF File path:",
                 fg_color='#b2baef', text_color='black', font=('Franklin Gothic Demi', 30, 'bold')).place(x=70, y=100)
    ctk.CTkEntry(New_Win, height=25, width=600, font=("helvetica", 20),
                 corner_radius=20, fg_color="white",
                 text_color="green", textvariable=file_path_var).place(x=60, y=145)
    ctk.CTkButton(New_Win, text="Browse", fg_color="white", text_color='green',
                  font=('Verdana', 20, 'bold'), width=12, height=12,
                  hover_color="#f9cb9c", border_width=2, border_color="grey",
                  command=open_existing_file).place(x=70, y=185)

    ctk.CTkLabel(New_Win, text="►Enter PDF to insert:", fg_color='#b2baef',
                 text_color='black', font=('Franklin Gothic Demi', 30, 'bold')).place(x=70, y=250)
    ctk.CTkEntry(New_Win, height=25, width=600, font=("helvetica", 20),
                 corner_radius=20, fg_color="white",
                 text_color="green", textvariable=insert_path_var).place(x=60, y=300)
    ctk.CTkButton(New_Win, text="Browse", fg_color="white", text_color='green',
                  hover_color="#f9cb9c", border_width=2, border_color="grey",
                  font=('Verdana', 20, 'bold'), width=12, height=12, command=open_insert_file).place(x=70, y=350)

    ctk.CTkLabel(New_Win, text="►Mention previous page number:", fg_color='#b2baef',
                 text_color='black', font=('Franklin Gothic Demi', 30, 'bold')).place(x=70, y=410)
    ctk.CTkEntry(New_Win, height=25, width=300, font=("helvetica", 20),
                 corner_radius=20, fg_color="white",
                 text_color="green", textvariable=position_var).place(x=70, y=460)
    ctk.CTkButton(New_Win, text="Insert Page", fg_color="white", text_color='red',
                  hover_color="#f9cb9c", border_width=5, border_color="black",
                  font=('Verdana', 25, 'bold'), width=15, height=15, command=insert_file).place(x=300, y=530)


def remove():
    New_Win = ctk.CTkToplevel(root)  # TopLevel object which will be treated as a new window
    New_Win.title("Edit PDF")
    New_Win.geometry('750x360')  # Order is WidthXHeight
    New_Win.minsize(750, 300)
    New_Win.maxsize(1000, 600)
    New_Win.configure(fg_color='#b2baef')

    heading = ctk.CTkLabel(New_Win, text="Remove page from an existing PDF ",
                           text_color="dark blue", fg_color="white", height=50,
                           font=("Microsoft New Tai Lue", 30, "bold"))
    heading.pack(fill='x', pady='15')

    file_path_var = StringVar()
    page_number_var = StringVar()

    def open_file():
        file = filedialog.askopenfilename(filetypes=[("PDF files", "*.pdf")])
        if file:
            file_path_var.set(file)

    def remove_page_from_p(input_pdf_path, page_number, output_pdf_path):
        try:
            # Read the PDF
            pdf_reader = PyPDF2.PdfReader(input_pdf_path)
            pdf_writer = PyPDF2.PdfWriter()

            # Remove the specified page
            for page_num in range(len(pdf_reader.pages)):
                if page_num != page_number - 1:  # Page numbers are 0-indexed 3==3 will not be included
                    pdf_writer.add_page(pdf_reader.pages[page_num])

            # Write the output PDF
            with open(output_pdf_path, 'wb') as output_pdf_file:
                pdf_writer.write(output_pdf_file)

            return output_pdf_path
        except Exception as e:
            print(f"Error removing page from PDF: {e}")
            return None

    def remove_page_from_pdf():
        input_pdf_path = file_path_var.get()
        page_number = int(page_number_var.get())
        output_pdf_path = filedialog.asksaveasfilename(defaultextension=".pdf", filetypes=[("PDF files", "*.pdf")])
        if output_pdf_path:
            result_path = remove_page_from_p(input_pdf_path, page_number, output_pdf_path)
            if result_path:
                messagebox.showinfo("Information", f"Page {page_number} removed successfully!")
                New_Win.destroy()
            else:
                messagebox.showerror("Error", "Failed to remove page.")
                New_Win.destroy()

    ctk.CTkLabel(New_Win, text="►Source PDF:",
                 fg_color='#b2baef', text_color='black', font=('Franklin Gothic Demi', 30, 'bold')).place(x=70, y=78)
    ctk.CTkEntry(New_Win, height=25, width=600, font=("helvetica", 20),
                 corner_radius=20, fg_color="white",
                 text_color="green", textvariable=file_path_var).place(x=70, y=120)
    ctk.CTkButton(New_Win, text="Browse", fg_color="white", text_color='green',
                  font=('Verdana', 20, 'bold'), width=12, height=12,
                  hover_color="#f9cb9c", border_width=2, border_color="grey", command=open_file).place(x=70, y=160)

    ctk.CTkLabel(New_Win, text="►Mention page number to be removed:",
                 fg_color='#b2baef', text_color='black', font=('Franklin Gothic Demi', 30, 'bold')).place(x=70, y=220)
    ctk.CTkEntry(New_Win, height=25, width=600, font=("helvetica", 20),
                 corner_radius=20, fg_color="white",
                 text_color="green", textvariable=page_number_var).place(x=70, y=265)
    ctk.CTkButton(New_Win, text="Remove", fg_color="white", text_color='red',
                  font=('Helvetica', 20, 'bold'), width=15, height=15,
                  hover_color="#f9cb9c", border_width=2, border_color="grey", command=remove_page_from_pdf).place(x=360,
                                                                                                                  y=310)


def edit_pdf():
    New_Win = ctk.CTkToplevel(root)  # TopLevel object which will be treated as a new window
    New_Win.title("Edit PDF")
    New_Win.geometry('640x340')  # Order is WidthXHeight
    New_Win.minsize(640, 300)
    New_Win.maxsize(700, 600)
    New_Win.configure(fg_color='#b2baef')

    heading = ctk.CTkLabel(New_Win, text="PDF Editing ",
                           text_color="dark blue", fg_color="white", height=50,
                           font=("Microsoft New Tai Lue", 30, "bold"))
    heading.pack(fill='x', pady='14')

    o = ctk.CTkButton(New_Win, fg_color="white", text="Insert pages\nto existing pdf",
                      command=insert, text_color="dark blue", hover_color="#f9cb9c",
                      font=("Arial Rounded MT Bold", 20, "bold"), image=img11, compound="top",
                      border_width=5, border_color="#291856")
    o.place(x=60, y=105)

    o = ctk.CTkButton(New_Win, fg_color="white", text="Remove pages from\nan existing pdf",
                      command=remove, text_color="dark blue", hover_color="#f9cb9c",
                      font=("Arial Rounded MT Bold", 20, "bold"), image=img12, compound="top",
                      border_width=5, border_color="#291856")
    o.place(x=360, y=105)


b1 = ctk.CTkButton(Frame2, text="Document to PDF", command=doc_to_pdf, text_color="Black",
                   font=("Arial Rounded MT Bold", 35, "bold"), fg_color="#FCE0C0", corner_radius=10,
                   border_width=5, border_spacing=10, border_color="brown", hover_color="#f9cb9c")
b1.place(x=180, y=110)

b2 = ctk.CTkButton(Frame2, text="PDF to Document", command=pdf_to_doc, text_color="Black",
                   font=("Arial Rounded MT Bold", 35, "bold"), fg_color="#FCE0C0", corner_radius=10,
                   border_width=5, border_spacing=10, border_color="brown", hover_color="#f9cb9c")
b2.place(x=180, y=200)

b3 = ctk.CTkButton(Frame2, text="Merge PDFs", command=merging, text_color="Black",
                   font=("Arial Rounded MT Bold", 35, "bold"), fg_color="#FCE0C0", corner_radius=10,
                   border_width=5, border_spacing=10, border_color="brown", hover_color="#f9cb9c")
b3.place(x=220, y=290)

b4 = ctk.CTkButton(Frame2, text="Edit PDF", command=edit_pdf,
                   font=("Arial Rounded MT Bold", 35, "bold"), text_color="Black",
                   fg_color="#FCE0C0", corner_radius=10,
                   border_width=5, border_spacing=10, border_color="brown", hover_color="#f9cb9c")
b4.place(x=240, y=380)

b5_close = ctk.CTkButton(Frame2, fg_color="#FCE0C0", text="Close", text_color="Black",
                         font=("Arial Rounded MT Bold", 20, "bold"),
                         image=img3, command=close,
                         border_width=5, border_spacing=10, border_color="black", hover_color="light grey")
b5_close.place(x=510, y=370)
# Initialize animation when the program starts


root.mainloop()
