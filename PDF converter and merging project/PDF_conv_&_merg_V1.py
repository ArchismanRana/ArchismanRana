from tkinter import *
from tkinter import messagebox, Toplevel
from PIL import ImageTk, Image, ImageDraw, ImageFont
from tkinter import filedialog
import pathlib
from pdf2docx import Converter
from win32com.client import Dispatch
import os
import sys
import PyPDF2
import img2pdf
import pdfplumber
import pandas as pd  # pip install openpyxl
from pdf2image import convert_from_path
from pptx import Presentation
from pptx.util import Inches
from pptx.enum.shapes import MSO_SHAPE_TYPE
import logging
from docx import Document
from docx.shared import Inches
from docx2pdf import convert


def resource_path(relative_path):
    """ Get absolute path to resource, works for dev and for PyInstaller """
    try:
        base_path = sys._MEIPASS
    except Exception:
        base_path = os.path.abspath(".")

    return os.path.join(base_path, relative_path)


root = Tk()
root.title("PDF converter and merging PDFs")

root.configure(background="#ca6b35")

# Always save icons photo in the file where this script will be saved
icon_path = resource_path('photos/Cavendish.png')
p = PhotoImage(file=icon_path)
root.iconphoto(True, p)
root.geometry('920x550')  # Order is WidthXHeight
root.minsize(600, 256)  # Order is Width,Height
root.configure(bg='#eab676')

Text_1 = Label(text="Welcome to PDF converter and Merging PDF", fg="dark blue", bg="white",
               font=("cosmeticians", 20, "bold"), borderwidth=7, relief=RIDGE)
# relief includes SUNKEN, RAISED, GROOVE, RIDGE
Text_1.pack(fill='x', pady='14')  # This .pack() is a rule in Python to show the stuffs

image1 = Image.open(resource_path('photos/WordToPDF.png'))
# Resize the Image  
image_1 = image1.resize((166, 100), Image.Resampling.LANCZOS)
# Convert the image to PhotoImage
img1 = ImageTk.PhotoImage(image_1)

image2 = Image.open(resource_path("photos/pdftoword.png"))
# Resize the Image
image_2 = image2.resize((166, 100), Image.Resampling.LANCZOS)
# Convert the image to PhotoImage
img2 = ImageTk.PhotoImage(image_2)

image3 = Image.open(resource_path("photos/Exit.png"))
# Resize the Image
image_3 = image3.resize((60, 55), Image.Resampling.LANCZOS)
# Convert the image to PhotoImage
img3 = ImageTk.PhotoImage(image_3)

image4 = Image.open(resource_path("photos/ExcelToPDF.png"))
# Resize the Image
image_4 = image4.resize((166, 100), Image.Resampling.LANCZOS)
# Convert the image to PhotoImage
img4 = ImageTk.PhotoImage(image_4)

image5 = Image.open(resource_path("photos/pdftoexcel.png"))
# Resize the Image
image_5 = image5.resize((166, 100), Image.Resampling.LANCZOS)
# Convert the image to PhotoImage
img5 = ImageTk.PhotoImage(image_5)

image7 = Image.open(resource_path("photos/PptToPdf.png"))
# Resize the Image
image_7 = image7.resize((166, 100), Image.Resampling.LANCZOS)
# Convert the image to PhotoImage
img7 = ImageTk.PhotoImage(image_7)

image8 = Image.open(resource_path("photos/Pdf_to_PPT.png"))
# Resize the Image
image_8 = image8.resize((166, 100), Image.Resampling.LANCZOS)
# Convert the image to PhotoImage
img8 = ImageTk.PhotoImage(image_8)

image9 = Image.open(resource_path("photos/JPGToPDF.png"))
# Resize the Image
image_9 = image9.resize((166, 100), Image.Resampling.LANCZOS)
# Convert the image to PhotoImage
img9 = ImageTk.PhotoImage(image_9)

image10 = Image.open(resource_path("photos/PdfToPng.png"))
# Resize the Image
image_10 = image10.resize((166, 100), Image.Resampling.LANCZOS)
# Convert the image to PhotoImage
img10 = ImageTk.PhotoImage(image_10)

image11 = Image.open(resource_path("photos/Insert_Page.png"))
# Resize the Image
image_11 = image11.resize((140, 100), Image.Resampling.LANCZOS)
# Convert the image to PhotoImage
img11 = ImageTk.PhotoImage(image_11)

image12 = Image.open(resource_path("photos/Remove_Page.png"))
# Resize the Image
image_12 = image12.resize((100, 100), Image.Resampling.LANCZOS)
# Convert the image to PhotoImage
img12 = ImageTk.PhotoImage(image_12)


def w_to_p():
    New_Win = Toplevel(root)  # TopLevel object which will be treated as a new window
    New_Win.title("Word to pdf")
    New_Win.geometry('800x350')  # Order is WidthXHeight
    New_Win.minsize(800, 400)
    New_Win.maxsize(1000, 600)
    New_Win.configure(bg='#c2c3ed')

    heading = Label(New_Win, text="Word to PDF converter",
                    fg="dark blue", bg="white",
                    font=("cosmeticians", 20, "bold"),
                    borderwidth=7, relief=RIDGE)
    heading.pack(fill='x', pady='14')

    Doc = Label(New_Win, text="Enter File path:", bg='#c2c3ed', fg='black', font='Verdana 20 bold')
    Doc.place(x=70, y=90)
    file_path_var = StringVar()
    En = Entry(New_Win, font=40, width=60, textvariable=file_path_var)
    En.place(x=70, y=150)

    def open_file():
        file = filedialog.askopenfilename(filetypes=[("Word files", "*.docx")])
        if file:
            file_path_var.set(file)

    def convert_file():
        file_path = file_path_var.get()
        if pathlib.Path(file_path).suffix == '.docx':
            convert(file_path)
            messagebox.showinfo("Information", "File converted successfully!")
            New_Win.destroy()
        else:
            messagebox.showerror("Error", "Please select a file with .docx extension")
            New_Win.destroy()

    browse_button = Button(New_Win, text="Browse", bg="white", fg='green', relief=RIDGE, font='Verdana 15 bold',
                           width=10, height=1, command=open_file)
    browse_button.place(x=80, y=210)

    convert_button = Button(New_Win, text="Convert", bg="white", fg='green', relief=RIDGE, font='Verdana 15 bold',
                            width=10, height=1, command=convert_file)
    convert_button.place(x=80, y=265)


def convert_excel_to_pdf(excel_path, pdf_path):
    """Converts an Excel file to a PDF file using win32com.client."""
    try:
        excel = Dispatch("Excel.Application")
        excel.Visible = False
        workbook = excel.Workbooks.Open(excel_path)
        workbook.ExportAsFixedFormat(0, pdf_path)
        workbook.Close(False)
        excel.Quit()
        logging.info(f"Converted {excel_path} to {pdf_path}")
    except Exception as e:
        logging.error(f"Failed to convert Excel to PDF: {e}")
        messagebox.showerror("Error", f"Failed to convert Excel to PDF: {e}")


def e_to_p():
    New_Win = Toplevel(root)  # TopLevel object which will be treated as a new window
    New_Win.title("Excel to PDF")
    New_Win.geometry('800x350')  # Order is WidthXHeight
    New_Win.minsize(800, 400)
    New_Win.maxsize(1000, 600)
    New_Win.configure(bg='#c2c3ed')

    heading = Label(New_Win, text="Excel to PDF Converter",
                    fg="dark blue", bg="white",
                    font=("cosmeticians", 20, "bold"),
                    borderwidth=7, relief=RIDGE)
    heading.pack(fill='x', pady='14')

    Doc = Label(New_Win, text="Enter File path:", bg='#c2c3ed', fg='black', font='Verdana 20 bold')
    Doc.place(x=70, y=90)
    file_path_var = StringVar()
    En = Entry(New_Win, font=40, width=60, textvariable=file_path_var)
    En.place(x=70, y=150)

    def open_file():
        file = filedialog.askopenfilename(filetypes=[("Excel files", "*.xlsx")])
        if file:
            file_path_var.set(file)

    def convert_file():
        file_path = file_path_var.get()
        if pathlib.Path(file_path).suffix == '.xlsx':
            pdf_path = filedialog.asksaveasfilename(defaultextension=".pdf", filetypes=[("PDF files", "*.pdf")])
            if pdf_path:
                convert_excel_to_pdf(file_path, pdf_path)
                messagebox.showinfo("Information", "File converted successfully!")
                New_Win.destroy()
        else:
            messagebox.showerror("Error", "Please select a file with .xlsx extension")
            New_Win.destroy()

    browse_button = Button(New_Win, text="Browse", bg="white", fg='green', relief=RIDGE, font='Verdana 15 bold',
                           width=10, height=1, command=open_file)
    browse_button.place(x=80, y=210)

    convert_button = Button(New_Win, text="Convert", bg="white", fg='green', relief=RIDGE, font='Verdana 15 bold',
                            width=10, height=1, command=convert_file)
    convert_button.place(x=80, y=265)


def add_poppler_to_path():
    # Specify the path to the Poppler bin directory
    poppler_path = r'C:\poppler\poppler-24.02.0\Library\bin'
    # Add Poppler bin directory to the PATH environment variable
    os.environ['PATH'] += os.pathsep + poppler_path


# Configure logging
logging.basicConfig(filename="ppt_to_pdf.log", level=logging.DEBUG, format="%(pastime)s - %(levelness)s - %(message)s")


def ppt_to_images(ppt_path, output_folder):
    if not os.path.exists(output_folder):
        os.makedirs(output_folder)

    prs = Presentation(ppt_path)
    image_files = []

    for i, slide in enumerate(prs.slides):
        slide_width = prs.slide_width  # Slide width in EMU (English Metric Units)
        slide_height = prs.slide_height  # Slide height in EMU

        slide_image_path = os.path.join(output_folder, f"slide_{i + 1}.png")
        image_files.append(slide_image_path)

        img = Image.new("RGB", (slide_width, slide_height), "white")
        draw = ImageDraw.Draw(img)

        for shape in slide.shapes:
            try:
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    image_path = os.path.join(output_folder, f"temp_image_{i}_{shape.shape_id}.png")
                    with open(image_path, 'wb') as f:
                        f.write(shape.image.blob)
                    image = Image.open(image_path)
                    img.paste(image, (int(shape.left), int(shape.top)))
                    os.remove(image_path)
                elif shape.has_text_frame:
                    text = shape.text_frame.text
                    font = ImageFont.load_default()
                    draw.text((int(shape.left), int(shape.top)), text, font=font, fill="black")
            except Exception as e:
                logging.error(f"Error processing shape on slide {i + 1}: {e}")
                continue

        img.save(slide_image_path)

    return image_files


def images_to_pdf(image_files, output_pdf_path):
    try:
        with open(output_pdf_path, "wb") as f:
            f.write(img2pdf.convert(image_files))
    except Exception as e:
        logging.error(f"Error converting images to PDF: {e}")


def ppt_to_word(ppt_path, word_path):
    prs = Presentation(ppt_path)
    doc = Document()

    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                q = doc.add_paragraph()
                q.add_run(shape.text).bold = True
            if shape.shape_type == 13:  # 13 corresponds to PICTURE
                image = shape.image
                image_path = f"{ppt_path}_{shape.shape_id}.png"
                with open(image_path, "wb") as f:
                    f.write(image.blob)
                doc.add_picture(image_path, width=Inches(6))
                os.remove(image_path)

    doc.save(word_path)
    return word_path


# Function to convert Word to PDF
def word_to_pdf(word_path, pdf_path):
    try:
        convert(word_path, pdf_path)
        return pdf_path
    except Exception as e:
        logging.error(f"Error converting Word to PDF: {e}")
        return None


# Function to convert PPT to PDF
def convert_ppt_to_pdf(ppt_path, pdf_path):
    try:
        word_path = ppt_path.replace('.pptx', '.docx').replace('.ppt', '.docx')
        ppt_to_word(ppt_path, word_path)
        pdf_output_path = word_to_pdf(word_path, pdf_path)
        if pdf_output_path:
            os.remove(word_path)  # Clean up the intermediate Word file
        return pdf_output_path
    except Exception as e:
        logging.error(f"Error converting PPT to PDF: {e}")
        return None


# Tkinter UI
def p_to_p():
    New_Win = Toplevel(root)  # TopLevel object which will be treated as a new window
    New_Win.title("PPT to PDF")
    New_Win.geometry('800x350')  # Order is WidthXHeight
    New_Win.minsize(800, 400)
    New_Win.maxsize(1000, 600)
    New_Win.configure(bg='#c2c3ed')

    heading = Label(New_Win, text="PPT to PDF Converter",
                    fg="dark blue", bg="white",
                    font=("cosmeticians", 20, "bold"),
                    borderwidth=7, relief=RIDGE)
    heading.pack(fill='x', pady='14')

    Label(New_Win, text="Enter File path:", bg='#c2c3ed', fg='black', font='Verdana 20 bold').place(x=70, y=90)
    file_path_var = StringVar()
    Entry(New_Win, font=40, width=60, textvariable=file_path_var).place(x=70, y=150)

    def open_file():
        file = filedialog.askopenfilename(filetypes=[("PowerPoint files", "*.ppt;*.pptx")])
        if file:
            file_path_var.set(file)

    def convert_file():
        file_path = file_path_var.get()
        if pathlib.Path(file_path).suffix in ['.ppt', '.pptx']:
            pdf_path = file_path.replace('.pptx', '.pdf').replace('.ppt', '.pdf')
            logging.info(f"Converting {file_path} to {pdf_path}")
            output_path = convert_ppt_to_pdf(file_path, pdf_path)
            if output_path:
                messagebox.showinfo("Information", f"File converted successfully! Saved at {output_path}")
                New_Win.destroy()
            else:
                messagebox.showerror("Error", "Failed to convert file.")
                New_Win.destroy()
        else:
            messagebox.showerror("Error", "Please select a file with .ppt or .pptx extension")

    Button(New_Win, text="Browse", bg="white", fg='green', relief=RIDGE, font='Verdana 15 bold',
           width=10, height=1, command=open_file).place(x=80, y=210)

    Button(New_Win, text="Convert", bg="white", fg='green', relief=RIDGE, font='Verdana 15 bold',
           width=10, height=1, command=convert_file).place(x=80, y=265)


def i_to_p():
    New_Win = Toplevel(root)  # TopLevel object which will be treated as a new window
    New_Win.title("Image to pdf")
    New_Win.geometry('800x350')  # Order is WidthXHeight
    New_Win.minsize(800, 400)
    New_Win.maxsize(1000, 600)
    New_Win.configure(bg='#c2c3ed')

    heading = Label(New_Win, text="Image to PDF converter",
                    fg="dark blue", bg="white",
                    font=("cosmeticians", 20, "bold"),
                    borderwidth=7, relief=RIDGE)
    heading.pack(fill='x', pady='14')

    Doc = Label(New_Win, text="Enter File path:", bg='#c2c3ed', fg='black', font='Verdana 20 bold')
    Doc.place(x=70, y=90)
    file_path_var = StringVar()
    En = Entry(New_Win, font=40, width=60, textvariable=file_path_var)
    En.place(x=70, y=150)

    def open_file():
        file = filedialog.askopenfilename(filetypes=[("Image Files", "*.png;*.jpg;*.jpeg;*.bmp;*.tiff")])
        if file:
            file_path_var.set(file)

    def convert_file():
        file_path = file_path_var.get()
        if pathlib.Path(file_path).suffix in ['.png', '.jpg', '.jpeg', '.bmp', '.tiff']:
            save_as = filedialog.asksaveasfilename(defaultextension=".pdf", filetypes=[("PDF Files", "*.pdf")])
            if save_as:
                try:
                    with open(save_as, "wb") as f:
                        f.write(img2pdf.convert(file_path))
                    messagebox.showinfo("Information", "File converted successfully!")
                    New_Win.destroy()
                except Exception as e:
                    messagebox.showerror("Error", f"Failed to convert file: {e}")
                    New_Win.destroy()
        else:
            messagebox.showerror("Error", "Please select a valid image file")
            New_Win.destroy()

    browse_button = Button(New_Win, text="Browse", bg="white", fg='green', relief=RIDGE, font='Verdana 15 bold',
                           width=10, height=1, command=open_file)
    browse_button.place(x=80, y=210)

    convert_button = Button(New_Win, text="Convert", bg="white", fg='green', relief=RIDGE, font='Verdana 15 bold',
                            width=10, height=1, command=convert_file)
    convert_button.place(x=80, y=265)


def merge_pdfs(filepaths, output_path):
    pdf_writer = PyPDF2.PdfWriter()  # Blank PDF
    for i in filepaths:
        if i:
            try:
                pdf_reader = PyPDF2.PdfReader(i)
                for page_num in range(len(pdf_reader.pages)):
                    page = pdf_reader.pages[page_num]
                    pdf_writer.add_page(page)  # PDFs added to blank pdf page by page
            except FileNotFoundError:
                return
            except Exception as e:
                messagebox.showerror("Error", f"An error occurred: {e}")
                return

        with open(output_path, 'wb') as out:
            pdf_writer.write(out)


def merging():
    New_Win = Toplevel(root)  # TopLevel object which will be treated as a new window
    New_Win.title("Merge pdf's")
    New_Win.geometry('860x770')  # Order is WidthXHeight
    New_Win.minsize(800, 400)
    New_Win.configure(bg='#c2c3ed')

    heading = Label(New_Win, text="Merge Four PDFs into a single PDF",
                    fg="dark blue", bg="white",
                    font=("cosmeticians", 20, "bold"),
                    borderwidth=7, relief=RIDGE)
    heading.pack(fill='x', pady='14')

    Doc_1 = Label(New_Win, text="Enter first File path:", bg='#c2c3ed', fg='black', font='Verdana 20 bold')
    Doc_1.place(x=70, y=100)
    file_path_var_1 = StringVar()
    En1 = Entry(New_Win, font=40, width=60, textvariable=file_path_var_1)
    En1.place(x=70, y=145)

    Doc_2 = Label(New_Win, text="Enter second File path:", bg='#c2c3ed', fg='black', font='Verdana 20 bold')
    Doc_2.place(x=70, y=250)
    file_path_var_2 = StringVar()
    En2 = Entry(New_Win, font=40, width=60, textvariable=file_path_var_2)
    En2.place(x=70, y=295)

    Doc_3 = Label(New_Win, text="Enter third File path:", bg='#c2c3ed', fg='black', font='Verdana 20 bold')
    Doc_3.place(x=70, y=400)
    file_path_var_3 = StringVar()
    En3 = Entry(New_Win, font=40, width=60, textvariable=file_path_var_3)
    En3.place(x=70, y=445)

    Doc_4 = Label(New_Win, text="Enter fourth File path:", bg='#c2c3ed', fg='black', font='Verdana 20 bold')
    Doc_4.place(x=70, y=550)
    file_path_var_4 = StringVar()
    En4 = Entry(New_Win, font=40, width=60, textvariable=file_path_var_4)
    En4.place(x=70, y=595)

    def open_1file():
        files = filedialog.askopenfilename(filetypes=[("PDF files", "*.pdf")])
        if files:
            file_path_var_1.set(files)

    def open_2file():
        files = filedialog.askopenfilename(filetypes=[("PDF files", "*.pdf")])
        if files:
            file_path_var_2.set(files)

    def open_3file():
        files = filedialog.askopenfilename(filetypes=[("PDF files", "*.pdf")])
        if files:
            file_path_var_3.set(files)

    def open_4file():
        files = filedialog.askopenfilename(filetypes=[("PDF files", "*.pdf")])
        if files:
            file_path_var_4.set(files)

    Doc_1entry = Button(New_Win, text="Browse File", bg="white", fg='green', relief=RIDGE,
                        font='Verdana 15 bold',
                        width=12,
                        height=1, command=open_1file)
    Doc_1entry.place(x=70, y=180)
    Doc_2entry = Button(New_Win, text="Browse File", bg="white", fg='green', relief=RIDGE,
                        font='Verdana 15 bold',
                        width=12,
                        height=1, command=open_2file)
    Doc_2entry.place(x=70, y=330)

    Doc_3entry = Button(New_Win, text="Browse File", bg="white", fg='green', relief=RIDGE,
                        font='Verdana 15 bold',
                        width=12,
                        height=1, command=open_3file)
    Doc_3entry.place(x=70, y=480)

    Doc_4entry = Button(New_Win, text="Browse File", bg="white", fg='green', relief=RIDGE,
                        font='Verdana 15 bold',
                        width=12,
                        height=1, command=open_4file)
    Doc_4entry.place(x=70, y=630)

    def merge_files():
        f1 = file_path_var_1.get()
        f2 = file_path_var_2.get()
        f3 = file_path_var_3.get()
        f4 = file_path_var_4.get()
        files = [f1, f2, f3, f4]
        valid_files = [f for f in files if f]  # Filter out empty file paths
        if len(valid_files) < 2:
            messagebox.showerror("Error", "Select at least two PDF files to merge.")
            return
        output_path = filedialog.asksaveasfilename(defaultextension=".pdf", filetypes=[("PDF files", "*.pdf")])
        if output_path:
            merge_pdfs(valid_files, output_path)
            messagebox.showinfo("Information", "PDFs merged successfully!")
            New_Win.destroy()

    merge_button = Button(New_Win, text="MERGE", bg="white", fg='green', relief=RIDGE, font='Verdana 15 bold',
                          width=10, height=1, command=merge_files)
    merge_button.place(x=600, y=630)


def close():
    x = messagebox.askquestion("askquestion", "Are you sure?")
    if x == 'yes':
        root.destroy()
    else:
        return None


def doc_to_pdf():
    New_Win = Toplevel(root)  # TopLevel object which will be treated as a new window
    New_Win.title("Word to pdf")
    New_Win.geometry('800x550')  # Order is WidthXHeight
    New_Win.minsize(800, 500)
    New_Win.maxsize(1000, 600)
    New_Win.configure(bg='#c2c3ed')

    heading = Label(New_Win, text="Convert Doc to PDF",
                    fg="dark blue", bg="white",
                    font=("cosmeticians", 20, "bold"),
                    borderwidth=7, relief=RIDGE)
    heading.pack(fill='x', pady='14')

    o = Button(New_Win, bg="white", text="Word to PDF", command=w_to_p, font=("Comic Sans", 15, "bold"),
               image=img1, borderwidth=5, relief=RIDGE, compound="top")
    o.place(x=100, y=105)

    o = Button(New_Win, bg="white", text="Excel to PDF", command=e_to_p, font=("Comic Sans", 15, "bold"),
               image=img4, borderwidth=5, relief=RIDGE, compound="top")
    o.place(x=500, y=105)

    o = Button(New_Win, bg="white", text="PPT to PDF", command=p_to_p, font=("Comic Sans", 15, "bold"),
               image=img7, borderwidth=5, relief=RIDGE, compound="top")
    o.place(x=100, y=305)

    o = Button(New_Win, bg="white", text="Image to PDF", command=i_to_p, font=("Comic Sans", 15, "bold"),
               image=img9, borderwidth=5, relief=RIDGE, compound="top")
    o.place(x=500, y=305)


def p_to_w():
    New_Win = Toplevel(root)  # TopLevel object which will be treated as a new window
    New_Win.title("pdf to Word")
    New_Win.geometry('800x350')  # Order is WidthXHeight
    New_Win.minsize(800, 400)
    New_Win.maxsize(1000, 600)
    New_Win.configure(bg='#c2c3ed')

    heading = Label(New_Win, text="PDF to Word converter",
                    fg="dark blue", bg="white",
                    font=("cosmeticians", 20, "bold"),
                    borderwidth=7, relief=RIDGE)
    heading.pack(fill='x', pady='14')

    Doc = Label(New_Win, text="Enter File path:", bg='#c2c3ed', fg='black', font='Verdana 20 bold')
    Doc.place(x=70, y=90)
    file_path_var = StringVar()
    En = Entry(New_Win, font=40, width=60, textvariable=file_path_var)
    En.place(x=70, y=150)

    def open_file():
        file = filedialog.askopenfilename(filetypes=[("Pdf files", "*.pdf")])
        if file:
            file_path_var.set(file)

    def convert_file():
        file_path = file_path_var.get()
        if pathlib.Path(file_path).suffix == '.pdf':
            docx_path = file_path.replace(".pdf", ".docx")
            cv = Converter(file_path)
            cv.convert(docx_path)
            cv.close()
            messagebox.showinfo("Information", "File converted successfully!")
            New_Win.destroy()
        else:
            messagebox.showerror("Error", "Failed to convert file. Please select a file with .docx extension")
            New_Win.destroy()

    browse_button = Button(New_Win, text="Browse", bg="white", fg='green', relief=RIDGE, font='Verdana 15 bold',
                           width=10, height=1, command=open_file)
    browse_button.place(x=80, y=210)

    convert_button = Button(New_Win, text="Convert", bg="white", fg='green', relief=RIDGE, font='Verdana 15 bold',
                            width=10, height=1, command=convert_file)
    convert_button.place(x=80, y=265)


def p_to_e():
    New_Win = Toplevel(root)  # TopLevel object which will be treated as a new window
    New_Win.title("PDF to Excel")
    New_Win.geometry('800x350')  # Order is WidthXHeight
    New_Win.minsize(800, 400)
    New_Win.maxsize(1000, 600)
    New_Win.configure(bg='#c2c3ed')

    heading = Label(New_Win, text="PDF to Excel converter",
                    fg="dark blue", bg="white",
                    font=("cosmeticians", 20, "bold"),
                    borderwidth=7, relief=RIDGE)
    heading.pack(fill='x', pady='14')

    Doc = Label(New_Win, text="Enter File path:", bg='#c2c3ed', fg='black', font='Verdana 20 bold')
    Doc.place(x=70, y=90)
    file_path_var = StringVar()
    En = Entry(New_Win, font=40, width=60, textvariable=file_path_var)
    En.place(x=70, y=150)

    def open_file():
        file = filedialog.askopenfilename(filetypes=[("Pdf files", "*.pdf")])
        if file:
            file_path_var.set(file)

    def convert_file():
        file_path = file_path_var.get()
        if pathlib.Path(file_path).suffix == '.pdf':
            save_as = filedialog.asksaveasfilename(defaultextension=".xlsx", filetypes=[("Excel Files", "*.xlsx")])
            if save_as:
                try:
                    with pdfplumber.open(file_path) as pdf:
                        all_tables = []
                        for page in pdf.pages:
                            tables = page.extract_tables()
                            for table in tables:
                                df = pd.DataFrame(table)
                                all_tables.append(df)

                        if all_tables:
                            result = pd.concat(all_tables, ignore_index=True)
                            result.to_excel(save_as, index=False)
                            messagebox.showinfo("Information", "File converted successfully!")
                            New_Win.destroy()
                        else:
                            messagebox.showerror("Error", "No tables found in the PDF")
                except Exception as e:
                    messagebox.showerror("Error", f"Failed to convert file: {e}")
                    New_Win.destroy()
        else:
            messagebox.showerror("Error", "Please select a valid PDF file")
            New_Win.destroy()

    browse_button = Button(New_Win, text="Browse", bg="white", fg='green', relief=RIDGE, font='Verdana 15 bold',
                           width=10, height=1, command=open_file)
    browse_button.place(x=80, y=210)

    convert_button = Button(New_Win, text="Convert", bg="white", fg='green', relief=RIDGE, font='Verdana 15 bold',
                            width=10, height=1, command=convert_file)
    convert_button.place(x=80, y=265)


def pdf_to_ppt(pdf_path, output_path):
    images = convert_from_path(pdf_path)
    presentation = Presentation()

    for image in images:
        slide_layout = presentation.slide_layouts[6]
        slide = presentation.slides.add_slide(slide_layout)
        imag_path = os.path.join(output_path, "slide_image.jpg")
        image.save(imag_path, 'JPEG')
        left = top = Inches(0)
        slide.shapes.add_picture(imag_path, left, top, width=Inches(10), height=Inches(7.5))

    pptx_output_path = os.path.join(output_path, "output_presentation.pptx")
    presentation.save(pptx_output_path)


def pd_to_pp():
    add_poppler_to_path()
    New_Win = Toplevel(root)  # TopLevel object which will be treated as a new window
    New_Win.title("pdf to PPT")
    New_Win.geometry('800x350')  # Order is WidthXHeight
    New_Win.minsize(800, 400)
    New_Win.maxsize(1000, 600)
    New_Win.configure(bg='#c2c3ed')

    heading = Label(New_Win, text="PDF to PPT converter",
                    fg="dark blue", bg="white",
                    font=("cosmeticians", 20, "bold"),
                    borderwidth=7, relief=RIDGE)
    heading.pack(fill='x', pady='14')

    Doc = Label(New_Win, text="Enter File path:", bg='#c2c3ed', fg='black', font='Verdana 20 bold')
    Doc.place(x=70, y=90)
    file_path_var = StringVar()
    En = Entry(New_Win, font=40, width=60, textvariable=file_path_var)
    En.place(x=70, y=150)

    def open_file():
        file = filedialog.askopenfilename(filetypes=[("Pdf files", "*.pdf")])
        if file:
            file_path_var.set(file)

    def convert_file():
        file_path = file_path_var.get()
        if pathlib.Path(file_path).suffix == '.pdf':
            output_dir = filedialog.askdirectory()
            if output_dir:
                pdf_to_ppt(file_path, output_dir)
                messagebox.showinfo("Information", "File converted successfully!")
                New_Win.destroy()

        else:
            messagebox.showerror("Error", "Failed to convert file. Select a file with .pdf extension")
            New_Win.destroy()

    browse_button = Button(New_Win, text="Browse", bg="white", fg='green', relief=RIDGE, font='Verdana 15 bold',
                           width=10, height=1, command=open_file)
    browse_button.place(x=80, y=210)

    convert_button = Button(New_Win, text="Convert", bg="white", fg='green', relief=RIDGE, font='Verdana 15 bold',
                            width=10, height=1, command=convert_file)
    convert_button.place(x=80, y=265)


def p_to_i():
    # Function to convert PDF to images using pdf2image
    add_poppler_to_path()
    New_Win = Toplevel(root)  # TopLevel object which will be treated as a new window
    New_Win.title("pdf to Word")
    New_Win.geometry('800x350')  # Order is WidthXHeight
    New_Win.minsize(800, 400)
    New_Win.maxsize(1000, 600)
    New_Win.configure(bg='#c2c3ed')

    heading = Label(New_Win, text="PDF to Word converter",
                    fg="dark blue", bg="white",
                    font=("cosmeticians", 20, "bold"),
                    borderwidth=7, relief=RIDGE)
    heading.pack(fill='x', pady='14')

    Doc = Label(New_Win, text="Enter File path:", bg='#c2c3ed', fg='black', font='Verdana 20 bold')
    Doc.place(x=70, y=90)
    file_path_var = StringVar()
    En = Entry(New_Win, font=40, width=60, textvariable=file_path_var)
    En.place(x=70, y=150)

    def open_file():
        file = filedialog.askopenfilename(filetypes=[("Pdf files", "*.pdf")])
        if file:
            file_path_var.set(file)

    def convert_file():
        file_path = file_path_var.get()
        if pathlib.Path(file_path).suffix == '.pdf':
            save_dir = filedialog.askdirectory()
            if save_dir:
                try:
                    images = convert_from_path(file_path)
                    for i, image in enumerate(images):
                        image.save(f"{save_dir}/page_{i + 1}.png", "PNG")
                    messagebox.showinfo("Information", "PDF converted to images successfully!")
                    New_Win.destroy()
                except Exception as e:
                    messagebox.showerror("Error", f"Failed to convert file: {e}")
                    New_Win.destroy()
        else:
            messagebox.showerror("Error", "Please select a valid PDF file")
            New_Win.destroy()

    browse_button = Button(New_Win, text="Browse", bg="white", fg='green', relief=RIDGE, font='Verdana 15 bold',
                           width=10, height=1, command=open_file)
    browse_button.place(x=80, y=210)

    convert_button = Button(New_Win, text="Convert", bg="white", fg='green', relief=RIDGE, font='Verdana 15 bold',
                            width=10, height=1, command=convert_file)
    convert_button.place(x=80, y=265)


def pdf_to_doc():
    New_Win = Toplevel(root)  # TopLevel object which will be treated as a new window
    New_Win.title("Word to pdf")
    New_Win.geometry('800x550')  # Order is WidthXHeight
    New_Win.minsize(800, 500)
    New_Win.maxsize(1000, 600)
    New_Win.configure(bg='#c2c3ed')

    heading = Label(New_Win, text="Convert PDF to Doc",
                    fg="dark blue", bg="white",
                    font=("cosmeticians", 20, "bold"),
                    borderwidth=7, relief=RIDGE)
    heading.pack(fill='x', pady='14')

    o = Button(New_Win, bg="white", text="PDF to Word", command=p_to_w,
               font=("Comic Sans", 15, "bold"),
               image=img2, borderwidth=5, relief=RIDGE, compound="top")
    o.place(x=100, y=105)

    o = Button(New_Win, bg="white", text="PDF to Excel", command=p_to_e, font=("Comic Sans", 15, "bold"),
               image=img5, borderwidth=5, relief=RIDGE, compound="top")
    o.place(x=500, y=105)

    o = Button(New_Win, bg="white", text="PDF to PPT", command=pd_to_pp, font=("Comic Sans", 15, "bold"),
               image=img8, borderwidth=5, relief=RIDGE, compound="top")
    o.place(x=100, y=305)

    o = Button(New_Win, bg="white", text="PDF to Image", command=p_to_i, font=("Comic Sans", 15, "bold"),
               image=img10, borderwidth=5, relief=RIDGE, compound="top")
    o.place(x=500, y=305)


def insert_page_into_pdf(existing_pdf_path, insert_pdf_path, position, output_path):
    existing_pdf_reader = PyPDF2.PdfReader(existing_pdf_path)
    insert_pdf_reader = PyPDF2.PdfReader(insert_pdf_path)
    pdf_writer = PyPDF2.PdfWriter()

    for page_num in range(position):
        pdf_writer.add_page(existing_pdf_reader.pages[page_num])

    for page in insert_pdf_reader.pages:
        pdf_writer.add_page(page)

    for page_num in range(position, len(existing_pdf_reader.pages)):
        pdf_writer.add_page(existing_pdf_reader.pages[page_num])

    with open(output_path, 'wb') as out_file:
        pdf_writer.write(out_file)


def insert():
    New_Win = Toplevel(root)  # TopLevel object which will be treated as a new window
    New_Win.title("Edit PDF")
    New_Win.geometry('750x650')  # Order is WidthXHeight
    New_Win.minsize(750, 600)
    New_Win.maxsize(1000, 600)
    New_Win.configure(bg='#c2c3ed')

    heading = Label(New_Win, text="Insert page to an existing PDF",
                    fg="dark blue", bg="white",
                    font=("cosmeticians", 20, "bold"),
                    borderwidth=7, relief=RIDGE)
    heading.pack(fill='x', pady='14')

    file_path_var = StringVar()
    insert_path_var = StringVar()
    position_var = StringVar()

    def open_existing_file():
        file = filedialog.askopenfilename(filetypes=[("PDF files", "*.pdf")])
        if file:
            file_path_var.set(file)

    def open_insert_file():
        file = filedialog.askopenfilename(filetypes=[("PDF files", "*.pdf")])
        if file:
            insert_path_var.set(file)

    def insert_file():
        existing_pdf_path = file_path_var.get()
        insert_pdf_path = insert_path_var.get()
        position = int(position_var.get())
        output_path = filedialog.asksaveasfilename(defaultextension=".pdf", filetypes=[("PDF files", "*.pdf")])
        if output_path:
            insert_page_into_pdf(existing_pdf_path, insert_pdf_path, position, output_path)
            messagebox.showinfo("Information", "Page inserted successfully!")
            New_Win.destroy()

    Label(New_Win, text="Enter existing PDF File path:", bg='#c2c3ed', fg='black',
          font='Verdana 20 bold').place(x=70, y=80)
    Entry(New_Win, font=40, width=40, textvariable=file_path_var).place(x=70, y=118)
    Button(New_Win, text="Browse", bg="white", fg='green', relief=RIDGE,
           font='Verdana 15 bold', width=10, height=1, command=open_existing_file).place(x=70, y=155)

    Label(New_Win, text="Enter PDF to insert:", bg='#c2c3ed', fg='black', font='Verdana 20 bold').place(x=70, y=220)
    Entry(New_Win, font=40, width=40, textvariable=insert_path_var).place(x=70, y=265)
    Button(New_Win, text="Browse", bg="white", fg='green', relief=RIDGE,
           font='Verdana 15 bold', width=10, height=1, command=open_insert_file).place(x=70, y=300)

    Label(New_Win, text="Mention previous page number:", bg='#c2c3ed', fg='black', font='Verdana 20 bold').place(x=70,
                                                                                                                 y=370)
    Entry(New_Win, font=40, width=20, textvariable=position_var).place(x=70, y=420)
    Button(New_Win, text="Insert Page", bg="white", fg='green', relief=RIDGE, font='Verdana 15 bold', width=15,
           height=1, command=insert_file).place(x=70, y=470)


def remove():
    New_Win = Toplevel(root)  # TopLevel object which will be treated as a new window
    New_Win.title("Edit PDF")
    New_Win.geometry('750x350')  # Order is WidthXHeight
    New_Win.minsize(750, 300)
    New_Win.maxsize(1000, 600)
    New_Win.configure(bg='#c2c3ed')

    heading = Label(New_Win, text="Remove page from an existing PDF ",
                    fg="dark blue", bg="white",
                    font=("cosmeticians", 20, "bold"),
                    borderwidth=7, relief=RIDGE)
    heading.pack(fill='x', pady='14')

    file_path_var = StringVar()
    page_number_var = StringVar()

    def open_file():
        file = filedialog.askopenfilename(filetypes=[("PDF files", "*.pdf")])
        if file:
            file_path_var.set(file)

    def remove_page_from_p(input_pdf_path, page_number, output_pdf_path):
        try:
            # Read the PDF
            pdf_reader = PyPDF2.PdfReader(input_pdf_path)
            pdf_writer = PyPDF2.PdfWriter()

            # Remove the specified page
            for page_num in range(len(pdf_reader.pages)):
                if page_num != page_number - 1:  # Page numbers are 0-indexed 3==3 will not be included
                    pdf_writer.add_page(pdf_reader.pages[page_num])

            # Write the output PDF
            with open(output_pdf_path, 'wb') as output_pdf_file:
                pdf_writer.write(output_pdf_file)

            return output_pdf_path
        except Exception as e:
            print(f"Error removing page from PDF: {e}")
            return None

    def remove_page_from_pdf():
        input_pdf_path = file_path_var.get()
        page_number = int(page_number_var.get())
        output_pdf_path = filedialog.asksaveasfilename(defaultextension=".pdf", filetypes=[("PDF files", "*.pdf")])
        if output_pdf_path:
            result_path = remove_page_from_p(input_pdf_path, page_number, output_pdf_path)
            if result_path:
                messagebox.showinfo("Information", f"Page {page_number} removed successfully!")
                New_Win.destroy()
            else:
                messagebox.showerror("Error", "Failed to remove page.")
                New_Win.destroy()

    Label(New_Win, text="Source PDF:", bg='#c2c3ed', fg='black', font='Verdana 20 bold').place(x=70, y=78)
    Entry(New_Win, font=40, width=40, textvariable=file_path_var).place(x=70, y=120)
    Button(New_Win, text="Browse", bg="white", fg='green', relief=RIDGE,
           font='Verdana 15 bold', width=10, height=1, command=open_file).place(x=70, y=160)

    Label(New_Win, text="Mention page number to be removed:", bg='#c2c3ed',
          fg='black', font='Verdana 20 bold').place(x=70, y=220)
    Entry(New_Win, font=40, width=20, textvariable=page_number_var).place(x=70, y=265)
    Button(New_Win, text="Remove", bg="white", fg='green', relief=RIDGE, font='Verdana 15 bold', width=15,
           height=1, command=remove_page_from_pdf).place(x=350, y=265)


def edit_pdf():
    New_Win = Toplevel(root)  # TopLevel object which will be treated as a new window
    New_Win.title("Edit PDF")
    New_Win.geometry('750x350')  # Order is WidthXHeight
    New_Win.minsize(750, 300)
    New_Win.maxsize(1000, 600)
    New_Win.configure(bg='#c2c3ed')

    heading = Label(New_Win, text="PDF Editing ",
                    fg="dark blue", bg="white",
                    font=("cosmeticians", 20, "bold"),
                    borderwidth=7, relief=RIDGE)
    heading.pack(fill='x', pady='14')

    o = Button(New_Win, bg="white", text="Insert pages to existing pdf", command=insert,
               font=("Comic Sans", 15, "bold"),
               image=img11, borderwidth=5, relief=RIDGE, compound="top")
    o.place(x=80, y=105)

    o = Button(New_Win, bg="white", text="Remove pages from\nan existing pdf", command=remove,
               font=("Comic Sans", 15, "bold"),
               image=img12, borderwidth=5, relief=RIDGE, compound="top")
    o.place(x=400, y=105)


b1 = Button(root, text="Document to PDF", command=doc_to_pdf,
            font=("Comic Sans", 20, "bold"),
            fg="Black", bg="white",
            activeforeground="Red",
            activebackground="white",
            borderwidth=10, relief=RIDGE)
b1.place(x=310, y=100)

b2 = Button(root, text="PDF to Document", command=pdf_to_doc,
            font=("Comic Sans", 20, "bold"),
            fg="Black", bg="white",
            activeforeground="Red",
            activebackground="white",
            borderwidth=10, relief=RIDGE)
b2.place(x=310, y=200)

b3 = Button(root, text="Merge PDFs", command=merging,
            font=("Comic Sans", 20, "bold"),
            fg="Black", bg="white",
            activeforeground="Red",
            activebackground="white",
            borderwidth=10, relief=RIDGE)
b3.place(x=350, y=300)

b4 = Button(root, text="Edit PDF", command=edit_pdf,
            font=("Comic Sans", 20, "bold"),
            fg="Black", bg="white",
            activeforeground="Red",
            activebackground="white",
            borderwidth=10, relief=RIDGE)
b4.place(x=370, y=400)

b5_close = Button(bg="white", text="Close", font=("Comic Sans", 15, "bold"),
                  image=img3,
                  command=close, borderwidth=5,
                  relief=RIDGE, compound="top")
b5_close.place(x=710, y=400)

root.mainloop()
